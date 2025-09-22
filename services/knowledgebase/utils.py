from io import BytesIO

import pytesseract
from docx import Document
from pdf2image import convert_from_bytes
from pptx import Presentation
from PyPDF2 import PdfReader


def extract_text_from_pdf(file_content: bytes):
    extracted_text = []
    reader = PdfReader(BytesIO(file_content))
    for page_number, page in enumerate(reader.pages, start=1):
        text = page.extract_text()
        if not text.strip():
            # If text extraction fails, use OCR
            images = convert_from_bytes(
                file_content, first_page=page_number, last_page=page_number
            )
            text = pytesseract.image_to_string(images[0])
        extracted_text.append({"page": page_number, "content": text})
    return extracted_text


def extract_text_from_docx(file_content: bytes):
    extracted_text = []
    document = Document(BytesIO(file_content))
    for i, paragraph in enumerate(document.paragraphs, start=1):
        extracted_text.append({"page": i, "content": paragraph.text})
    return extracted_text


def extract_text_from_pptx(file_content: bytes):
    extracted_text = []
    presentation = Presentation(BytesIO(file_content))
    for slide_number, slide in enumerate(presentation.slides, start=1):
        slide_text = "\n".join(
            [shape.text for shape in slide.shapes if hasattr(shape, "text")]
        )
        extracted_text.append({"page": slide_number, "content": slide_text})
    return extracted_text


def extract_text_from_txt(file_content: bytes):
    return [{"page": 1, "content": file_content.decode("utf-8")}]
